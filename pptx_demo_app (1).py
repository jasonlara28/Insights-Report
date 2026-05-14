"""
Category Insights Dashboard → PowerPoint Exporter
====================================================
A Streamlit app that analyzes category sales data, displays interactive
Plotly charts, and exports a polished PowerPoint deck with one click.

Requirements:
    pip install streamlit plotly pandas python-pptx kaleido

Run:
    streamlit run pptx_demo_app.py
"""

import io
import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ─────────────────────────────────────────────
# CONFIG & CONSTANTS
# ─────────────────────────────────────────────
COLOR_PALETTE = ["#4E79A7", "#F28E2B", "#59A14F", "#E15759",
                 "#76B7B2", "#EDC948", "#B07AA1", "#FF9DA7"]

PLOTLY_LAYOUT = dict(
    font=dict(family="Segoe UI, Arial, sans-serif", size=13),
    plot_bgcolor="white",
    paper_bgcolor="white",
    margin=dict(l=60, r=30, t=50, b=50),
    legend=dict(orientation="h", yanchor="bottom", y=-0.25, xanchor="center", x=0.5),
)

MONTHS_ORDER = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]


# ─────────────────────────────────────────────
# SAMPLE DATA GENERATOR
# ─────────────────────────────────────────────
def generate_sample_data() -> pd.DataFrame:
    """Return 12 months × 4 categories of realistic CPG sales data."""
    import random
    random.seed(42)

    categories = {
        "Baking Chips":     {"base_sales": 320000, "base_units": 48000, "yoy": 4.2},
        "Cookie Dough":     {"base_sales": 275000, "base_units": 39000, "yoy": 7.8},
        "Evaporated Milk":  {"base_sales": 190000, "base_units": 55000, "yoy": -1.5},
        "Pumpkin Puree":    {"base_sales": 140000, "base_units": 31000, "yoy": 12.3},
    }

    # Seasonal multipliers (index 0 = Jan … 11 = Dec)
    seasonal = [0.80, 0.75, 0.85, 0.90, 0.95, 1.00,
                1.00, 1.05, 1.15, 1.30, 1.45, 1.50]

    rows = []
    for cat, params in categories.items():
        for i, month in enumerate(MONTHS_ORDER):
            noise = random.uniform(0.93, 1.07)
            sales = round(params["base_sales"] * seasonal[i] * noise, 2)
            units = round(params["base_units"] * seasonal[i] * noise)
            yoy = round(params["yoy"] + random.uniform(-2.0, 2.0), 1)
            rows.append({"Category": cat, "Month": month,
                         "Sales": sales, "Units": units, "YoY_Change": yoy})

    df = pd.DataFrame(rows)
    df["Month"] = pd.Categorical(df["Month"], categories=MONTHS_ORDER, ordered=True)
    return df.sort_values(["Category", "Month"]).reset_index(drop=True)


# ─────────────────────────────────────────────
# CHART BUILDERS
# ─────────────────────────────────────────────
def build_bar_chart(df: pd.DataFrame) -> go.Figure:
    """Grouped bar chart — Sales by Category by Month."""
    fig = px.bar(
        df, x="Month", y="Sales", color="Category",
        barmode="group", color_discrete_sequence=COLOR_PALETTE,
        title="Monthly Sales by Category",
        labels={"Sales": "Sales ($)", "Month": ""},
    )
    fig.update_layout(**PLOTLY_LAYOUT)
    fig.update_yaxes(gridcolor="#ECECEC", tickprefix="$", tickformat=",")
    fig.update_xaxes(showgrid=False)
    return fig


def build_line_chart(df: pd.DataFrame) -> go.Figure:
    """Line chart — Units trend over time by Category."""
    fig = px.line(
        df, x="Month", y="Units", color="Category",
        markers=True, color_discrete_sequence=COLOR_PALETTE,
        title="Unit Volume Trend by Category",
        labels={"Units": "Units Sold", "Month": ""},
    )
    fig.update_layout(**PLOTLY_LAYOUT)
    fig.update_yaxes(gridcolor="#ECECEC", tickformat=",")
    fig.update_xaxes(showgrid=False)
    fig.update_traces(line=dict(width=2.5), marker=dict(size=6))
    return fig


def build_donut_chart(df: pd.DataFrame) -> go.Figure:
    """Donut chart — Sales share by Category."""
    cat_totals = df.groupby("Category", observed=True)["Sales"].sum().reset_index()
    fig = px.pie(
        cat_totals, names="Category", values="Sales", hole=0.45,
        color_discrete_sequence=COLOR_PALETTE,
        title="Sales Share by Category",
    )
    fig.update_layout(**PLOTLY_LAYOUT)
    fig.update_traces(textinfo="label+percent", textfont_size=13)
    return fig


# ─────────────────────────────────────────────
# INSIGHTS GENERATOR
# ─────────────────────────────────────────────
def generate_insights(df: pd.DataFrame) -> list[str]:
    """Auto-generate bullet-point insights from the data."""
    cat_sales = df.groupby("Category", observed=True)["Sales"].sum()
    cat_units = df.groupby("Category", observed=True)["Units"].sum()
    cat_yoy   = df.groupby("Category", observed=True)["YoY_Change"].mean()

    top_sales_cat = cat_sales.idxmax()
    top_units_cat = cat_units.idxmax()
    fastest_grow  = cat_yoy.idxmax()
    slowest_grow  = cat_yoy.idxmin()

    insights = [
        f"Top revenue category: {top_sales_cat} (${cat_sales[top_sales_cat]:,.0f} total sales).",
        f"Highest unit volume: {top_units_cat} ({cat_units[top_units_cat]:,} units sold).",
        f"Fastest YoY growth: {fastest_grow} ({cat_yoy[fastest_grow]:+.1f}% avg monthly change).",
        f"Slowest YoY growth: {slowest_grow} ({cat_yoy[slowest_grow]:+.1f}% avg monthly change).",
        f"Peak sales month across all categories: {df.groupby('Month', observed=True)['Sales'].sum().idxmax()}.",
        f"Total annual revenue: ${df['Sales'].sum():,.0f} across {df['Category'].nunique()} categories.",
    ]
    return insights


# ─────────────────────────────────────────────
# POWERPOINT BUILDER
# ─────────────────────────────────────────────
def fig_to_png_bytes(fig: go.Figure, width=960, height=540) -> bytes:
    """Export a Plotly figure to high-resolution PNG bytes using Kaleido."""
    return fig.to_image(format="png", width=width, height=height, scale=2, engine="kaleido")


def build_pptx(charts: list[tuple[str, go.Figure]], insights: list[str]) -> bytes:
    """
    Build a PowerPoint deck and return it as bytes.

    Parameters
    ----------
    charts : list of (title, plotly_figure) tuples
    insights : list of insight strings
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    # ── Helpers ──────────────────────────────
    def _set_bg(slide, r, g, b):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    def _add_textbox(slide, left, top, width, height, text,
                     font_size=18, bold=False, color=(51, 51, 51), alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = alignment
        return tf

    # ── Slide 1 — Title ─────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_bg(slide, 78, 121, 167)  # Brand blue

    _add_textbox(slide,
                 left=Inches(1), top=Inches(2.2),
                 width=Inches(11.3), height=Inches(1.5),
                 text="Category Insights Report",
                 font_size=40, bold=True, color=(255, 255, 255),
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 left=Inches(1), top=Inches(3.6),
                 width=Inches(11.3), height=Inches(0.8),
                 text="Auto-generated dashboard  •  Powered by Streamlit + Plotly + python-pptx",
                 font_size=18, bold=False, color=(220, 230, 241),
                 alignment=PP_ALIGN.CENTER)

    # ── Chart Slides ─────────────────────────
    for title, fig in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, 255, 255, 255)

        # Slide title
        _add_textbox(slide,
                     left=Inches(0.6), top=Inches(0.3),
                     width=Inches(12), height=Inches(0.6),
                     text=title, font_size=24, bold=True)

        # Export chart as PNG and insert
        img_bytes = fig_to_png_bytes(fig)
        img_stream = io.BytesIO(img_bytes)

        img_width  = Inches(10.5)
        img_height = Inches(5.6)
        left = int((slide_width - img_width) / 2)
        top  = Inches(1.2)
        slide.shapes.add_picture(img_stream, left, top, img_width, img_height)

    # ── Insights Slide ───────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, 248, 249, 250)

    _add_textbox(slide,
                 left=Inches(0.6), top=Inches(0.3),
                 width=Inches(12), height=Inches(0.6),
                 text="Key Insights", font_size=28, bold=True)

    # Add each insight as a separate paragraph
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(11), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, insight in enumerate(insights):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {insight}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(14)

    # ── Save to bytes ────────────────────────
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ─────────────────────────────────────────────
# STREAMLIT APP
# ─────────────────────────────────────────────
def main():
    st.set_page_config(page_title="Category Insights → PPTX", layout="wide", page_icon="📊")

    st.title("📊 Category Insights Dashboard")
    st.caption("Analyze category data  →  Preview charts  →  Export to PowerPoint in one click")

    # ── Sidebar ──────────────────────────────
    st.sidebar.header("⚙️ Data Source")
    upload = st.sidebar.file_uploader(
        "Upload a CSV file", type=["csv"],
        help="Expected columns: Category, Month, Sales, Units, YoY_Change"
    )

    if upload is not None:
        df = pd.read_csv(upload)
        # Ensure Month ordering
        if "Month" in df.columns:
            df["Month"] = pd.Categorical(df["Month"], categories=MONTHS_ORDER, ordered=True)
    else:
        df = generate_sample_data()
        st.sidebar.info("Using built-in sample data. Upload a CSV to use your own.")

    # Category filter
    st.sidebar.header("🔎 Filters")
    all_cats = sorted(df["Category"].unique())
    selected = st.sidebar.multiselect("Select Categories", all_cats, default=all_cats)

    if not selected:
        st.warning("Please select at least one category.")
        st.stop()

    df_filt = df[df["Category"].isin(selected)].copy()

    # ── Metrics Row ──────────────────────────
    col1, col2, col3 = st.columns(3)
    col1.metric("💰 Total Sales",   f"${df_filt['Sales'].sum():,.0f}")
    col2.metric("📦 Total Units",   f"{df_filt['Units'].sum():,}")
    col3.metric("📈 Avg YoY Change", f"{df_filt['YoY_Change'].mean():+.1f}%")

    st.divider()

    # ── Charts ───────────────────────────────
    fig_bar   = build_bar_chart(df_filt)
    fig_line  = build_line_chart(df_filt)
    fig_donut = build_donut_chart(df_filt)

    # Two-column layout for bar and line; donut below
    left_col, right_col = st.columns(2)
    with left_col:
        st.plotly_chart(fig_bar, use_container_width=True)
    with right_col:
        st.plotly_chart(fig_line, use_container_width=True)

    st.plotly_chart(fig_donut, use_container_width=True)

    # ── Insights ─────────────────────────────
    st.divider()
    insights = generate_insights(df_filt)
    st.subheader("💡 Key Insights")
    for insight in insights:
        st.markdown(f"- {insight}")

    # ── PowerPoint Export ────────────────────
    st.divider()
    st.subheader("📥 Export to PowerPoint")
    st.write("Click below to generate a `.pptx` deck with all charts and insights.")

    charts_for_export = [
        ("Monthly Sales by Category",   fig_bar),
        ("Unit Volume Trend by Category", fig_line),
        ("Sales Share by Category",      fig_donut),
    ]

    if st.button("🔨 Generate PowerPoint", type="primary"):
        with st.spinner("Building your deck — exporting charts & assembling slides…"):
            pptx_bytes = build_pptx(charts_for_export, insights)

        st.success("✅ PowerPoint ready!")
        st.download_button(
            label="⬇️ Download PowerPoint",
            data=pptx_bytes,
            file_name="category_insights_report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


if __name__ == "__main__":
    main()
